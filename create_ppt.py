import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def create_ppt_with_images(output_path, image_data_list):
    """
    创建PPT，每个图片生成两页：一页文字描述，一页插入图片

    Args:
        output_path: 输出PPT文件的路径
        image_data_list: 列表，每个元素是字典，包含：
            - 'image_path': 图片路径
            - 'description': 文字描述
            - 'title': (可选) 标题
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, data in enumerate(image_data_list):
        image_path = data['image_path']
        description = data['description']
        title = data.get('title', f'第 {idx + 1} 页')

        if not os.path.exists(image_path):
            print(f"警告：图片不存在 - {image_path}")
            continue

        slide_layout_text = prs.slide_layouts[6]
        slide_text = prs.slides.add_slide(slide_layout_text)

        title_box = slide_text.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
        )
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER

        desc_box = slide_text.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.7)
        )
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True

        lines = description.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf_desc.paragraphs[0]
            else:
                p = tf_desc.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)

        slide_layout_image = prs.slide_layouts[6]
        slide_image = prs.slides.add_slide(slide_layout_image)

        img_title_box = slide_image.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        img_tf_title = img_title_box.text_frame
        img_p_title = img_tf_title.paragraphs[0]
        img_p_title.text = f'{title} - 图片展示'
        img_p_title.font.size = Pt(28)
        img_p_title.font.bold = True
        img_p_title.alignment = PP_ALIGN.CENTER

        try:
            pic = slide_image.shapes.add_picture(
                image_path,
                left=Inches(1.5),
                top=Inches(1.3),
                width=Inches(10.333),
                height=Inches(5.7)
            )
        except Exception as e:
            print(f"插入图片失败 {image_path}: {e}")

        print(f"已处理: {title} (文字页 + 图片页)")

    prs.save(output_path)
    print(f"\nPPT 已成功保存至: {output_path}")


def main():
    """
    主函数：从config.json读取配置数据并生成PPT
    使用方法：
    1. 编辑 config.json 文件，添加你的图片和描述
    2. 运行此脚本即可生成 PPT
    """

    config_file = "config.json"

    if not os.path.exists(config_file):
        print(f"错误：配置文件 {config_file} 不存在！")
        print("请先创建 config.json 文件")
        return

    try:
        with open(config_file, 'r', encoding='utf-8') as f:
            config = json.load(f)
    except json.JSONDecodeError as e:
        print(f"错误：JSON 格式错误 - {e}")
        return
    except Exception as e:
        print(f"错误：读取配置文件失败 - {e}")
        return

    output_file = config.get('output_file', 'output.pptx')
    image_data = config.get('images', [])

    if not image_data:
        print("警告：配置文件中没有图片数据！")
        return

    print("=" * 50)
    print("开始生成 PPT...")
    print(f"共 {len(image_data)} 组图片数据")
    print("=" * 50 + "\n")

    create_ppt_with_images(output_file, image_data)


if __name__ == "__main__":
    main()
